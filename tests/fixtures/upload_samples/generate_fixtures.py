"""Generate test fixtures for v7.5.0 upload-resilience tests.

Run: python tests/fixtures/upload_samples/generate_fixtures.py

Creates 12 sample files used by:
- scripts/upload_resilience_e2e_verify.py (Backend E2E)
- tests/test_extraction_v750.py (pytest unit)
- tests/e2e-ui/v7.5.0-upload-resilience.spec.js (Playwright)

Files generated are NOT committed to git (.gitignore-d) — regenerate before tests.
"""
from __future__ import annotations

import io
import json
import os
import sys
from pathlib import Path

OUT_DIR = Path(__file__).parent

# ─── helpers ────────────────────────────────────────────────────────


def gen_png_with_text():
    """Generate a 400x100 PNG with Thai+EN text using PIL."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (400, 100), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except OSError:
        font = ImageFont.load_default()
    draw.text((10, 30), "Hello PDB v7.5.0", fill="black", font=font)
    img.save(OUT_DIR / "sample.png", "PNG")
    print("  ✓ sample.png (PNG with text)")


def gen_blank_png():
    from PIL import Image
    img = Image.new("RGB", (400, 100), color="white")
    img.save(OUT_DIR / "sample_blank.png", "PNG")
    print("  ✓ sample_blank.png (blank, no text)")


def gen_xss_png():
    """PNG with XSS-style text content (alert payload as image text)."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (400, 100), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 20)
    except OSError:
        font = ImageFont.load_default()
    draw.text((10, 30), "alert(1) script test", fill="black", font=font)
    img.save(OUT_DIR / "sample_xss.png", "PNG")
    print("  ✓ sample_xss.png (PNG with XSS-style text)")


def gen_simple_pdf():
    """Generate simple PDF with known text (~1KB)."""
    try:
        from PyPDF2 import PdfWriter
        from PyPDF2.generic import RectangleObject
        # PyPDF2 alone can't easily create with text — use reportlab
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(OUT_DIR / "sample.pdf"))
        c.drawString(100, 750, "Hello PDB v7.5.0 — sample PDF for testing.")
        c.drawString(100, 730, "This is line 2 with longer content for extraction tests.")
        c.save()
        print("  ✓ sample.pdf (simple PDF)")
    except ImportError:
        # Fallback: write minimal valid PDF skeleton
        minimal_pdf = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 595 842]/Contents 4 0 R/Resources<</Font<</F1<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>>>>>>>endobj
4 0 obj<</Length 60>>stream
BT /F1 12 Tf 100 750 Td (Hello PDB v7.5.0 sample PDF) Tj ET
endstream endobj
xref
0 5
0000000000 65535 f
0000000009 00000 n
0000000054 00000 n
0000000099 00000 n
0000000220 00000 n
trailer<</Size 5/Root 1 0 R>>
startxref
320
%%EOF
"""
        (OUT_DIR / "sample.pdf").write_bytes(minimal_pdf)
        print("  ✓ sample.pdf (minimal PDF — no reportlab)")


def gen_empty_pdf():
    """0-byte file with .pdf extension (tests EMPTY_FILE)."""
    (OUT_DIR / "empty.pdf").write_bytes(b"")
    print("  ✓ empty.pdf (0 bytes)")


def gen_big_text_pdf():
    """Generate a PDF with ~150K chars (triggers Phase 4 chunking)."""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        c = canvas.Canvas(str(OUT_DIR / "big_text.pdf"), pagesize=A4)
        # Lorem ipsum repeated to reach ~150K chars
        para = (
            "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
            "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. "
            "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris. "
        )
        # 50 pages × ~3000 chars = 150K chars
        for page_n in range(1, 51):
            y = 800
            c.drawString(50, y, f"=== PAGE {page_n} — UNIQUE_MARKER_{page_n:03d} ===")
            y -= 20
            for line_n in range(35):
                c.drawString(50, y, f"L{line_n}: {para}")
                y -= 18
                if y < 50:
                    break
            c.showPage()
        c.save()
        size = (OUT_DIR / "big_text.pdf").stat().st_size
        print(f"  ✓ big_text.pdf ({size//1024}KB, 50 pages with markers)")
    except ImportError:
        # Fallback: just create big TXT file as substitute
        para = "Lorem ipsum dolor sit amet, consectetur adipiscing elit. " * 50
        content = "\n\n".join(
            f"# PAGE {n} — UNIQUE_MARKER_{n:03d}\n{para}"
            for n in range(1, 51)
        )
        (OUT_DIR / "big_text.txt").write_text(content, encoding="utf-8")
        print(f"  ⚠ big_text.txt (no reportlab — using TXT instead, {len(content)} chars)")


def gen_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    # Sheet 1
    ws1 = wb.active
    ws1.title = "Q1_Sales"
    ws1["A1"] = "Product"
    ws1["B1"] = "Revenue"
    ws1["A2"] = "Widget A"
    ws1["B2"] = 1000
    ws1["A3"] = "Widget B"
    ws1["B3"] = 2500
    ws1["B4"] = "=SUM(B2:B3)"  # Formula — should resolve to 3500 with data_only
    # Sheet 2
    ws2 = wb.create_sheet("Q2_Sales")
    ws2["A1"] = "Region"
    ws2["B1"] = "Total"
    ws2["A2"] = "North"
    ws2["B2"] = 5000
    # Sheet 3
    ws3 = wb.create_sheet("Notes")
    ws3["A1"] = "Important: confidential data"
    wb.save(OUT_DIR / "sample.xlsx")
    print("  ✓ sample.xlsx (3 sheets + formula)")


def gen_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "PDB v7.5.0 Test Deck"
    s1.placeholders[1].text = "Slide content for testing extraction"
    s1.notes_slide.notes_text_frame.text = "SPEAKER_NOTES_MARKER_S1"
    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Slide 2 Topic"
    s2.placeholders[1].text = "More content here"
    s2.notes_slide.notes_text_frame.text = "SPEAKER_NOTES_MARKER_S2"
    prs.save(OUT_DIR / "sample.pptx")
    print("  ✓ sample.pptx (2 slides + speaker notes)")


def gen_html_safe():
    html = """<!DOCTYPE html>
<html>
<head><title>Safe HTML Test</title></head>
<body>
<h1>PDB v7.5.0 HTML Extract Test</h1>
<p>This is a paragraph with <strong>bold text</strong>.</p>
<h2>Subsection</h2>
<p>Another paragraph for testing.</p>
</body>
</html>"""
    (OUT_DIR / "sample_safe.html").write_text(html, encoding="utf-8")
    print("  ✓ sample_safe.html (clean HTML)")


def gen_html_xss():
    html = """<!DOCTYPE html>
<html>
<head>
<title>XSS Test</title>
<style>body { color: red; } /* CSS_MARKER_SHOULD_BE_STRIPPED */</style>
<script>alert(1); /* JS_MARKER_SHOULD_BE_STRIPPED */</script>
</head>
<body>
<h1>Safe Content</h1>
<p>This text should be in extracted output.</p>
<script>console.log('inline_script_marker');</script>
</body>
</html>"""
    (OUT_DIR / "sample_xss.html").write_text(html, encoding="utf-8")
    print("  ✓ sample_xss.html (with <script> + <style> for security test)")


def gen_json():
    data = {
        "version": "7.5.0",
        "feature": "upload_resilience",
        "phases": [
            {"id": 1, "name": "Fix Bugs"},
            {"id": 4, "name": "Big File Support"},
            {"id": 2, "name": "Proactive UX"},
            {"id": 3, "name": "More Formats"},
        ],
        "settings": {"chunk_size": 10000, "overlap": 500},
    }
    (OUT_DIR / "sample.json").write_text(
        json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    print("  ✓ sample.json (nested structure)")


def gen_rtf():
    # Minimal RTF document
    rtf = (
        r"{\rtf1\ansi\deff0"
        r"{\fonttbl{\f0\froman Times New Roman;}}"
        r"\f0\fs24 "
        r"PDB v7.5.0 RTF test content. "
        r"This is a sample RTF file for extraction testing. "
        r"Bold word: \b boldword\b0 . End of file."
        r"}"
    )
    (OUT_DIR / "sample.rtf").write_bytes(rtf.encode("ascii"))
    print("  ✓ sample.rtf (basic RTF)")


def gen_unsupported():
    """Generate .xyz file (test UNSUPPORTED_TYPE)."""
    (OUT_DIR / "sample.xyz").write_text("This is unsupported", encoding="utf-8")
    print("  ✓ sample.xyz (unsupported extension)")


def gen_empty_txt():
    (OUT_DIR / "empty.txt").write_bytes(b"")
    print("  ✓ empty.txt (0 bytes)")


# ─── main ────────────────────────────────────────────────────────────


def main():
    print(f"Generating fixtures into {OUT_DIR}")

    generators = [
        ("PIL (image)", gen_png_with_text),
        ("PIL blank", gen_blank_png),
        ("PIL XSS image", gen_xss_png),
        ("Simple PDF", gen_simple_pdf),
        ("Empty PDF", gen_empty_pdf),
        ("Big text PDF", gen_big_text_pdf),
        ("XLSX", gen_xlsx),
        ("PPTX", gen_pptx),
        ("HTML safe", gen_html_safe),
        ("HTML XSS", gen_html_xss),
        ("JSON", gen_json),
        ("RTF", gen_rtf),
        ("Unsupported", gen_unsupported),
        ("Empty TXT", gen_empty_txt),
    ]

    failed = []
    for name, func in generators:
        try:
            func()
        except Exception as e:
            failed.append((name, str(e)))
            print(f"  ✗ {name} failed: {e}")

    print(f"\nDone — {len(generators) - len(failed)}/{len(generators)} generated")
    if failed:
        print(f"Failures: {failed}")
        sys.exit(1)


if __name__ == "__main__":
    main()
