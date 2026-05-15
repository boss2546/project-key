"""Local Office extractors -- v10.0.2 (HANDOFF Decision 4 + Pattern A/D).

Why local extract instead of LlamaParse / Docling for Office files:
  - Lab-validated 40-76x cheaper per file (HANDOFF section 3 quality table).
  - 5-6x faster end-to-end (15s vs 90s for 90KB DOCX).
  - LlamaParse drops 51-98% of Thai diacritics on DOCX font-encoded text.
  - python-docx / python-pptx / openpyxl read XML directly, no rendering.

Routing decision flow (HANDOFF Pattern A):
  1. ext in {docx, pptx, xlsx} + USE_LOCAL_EXTRACT_X=true + size <= MAX_MB → local
  2. else → caller falls back to LlamaParse (Pattern B)

Concurrency safety (HANDOFF Pattern D + E):
  - PARSE_LOCK (threading.Lock) wraps python-docx / python-pptx calls -- both
    use lxml's module-level oxml_parser which is NOT thread-safe under load.
    Stress test 100 concurrent files without lock → random XMLSyntaxError.
    openpyxl does NOT need the lock.

Output schema (HANDOFF Pattern G):
  Each function returns (text, warnings) tuple. warnings list is empty on
  the happy path -- non-fatal issues (large file, encoding fallback, empty
  sheet, etc.) populate it and surface in extract_warnings column.
"""
from __future__ import annotations

import logging
import os
from typing import Tuple, List

from ..config import (
    LOCAL_EXTRACT_MAX_MB,
    USE_LOCAL_EXTRACT_DOCX,
    USE_LOCAL_EXTRACT_PPTX,
    USE_LOCAL_EXTRACT_XLSX,
)
from .safety import PARSE_LOCK

logger = logging.getLogger(__name__)


__all__ = [
    "should_use_local",
    "read_docx_as_markdown",
    "read_pptx_as_markdown",
    "read_xlsx_as_markdown",
]


# ============================================================
# Routing helper
# ============================================================

def should_use_local(ext: str, file_size_bytes: int = 0) -> bool:
    """HANDOFF Pattern A — flag-based routing with size threshold.

    Returns False if either:
      - feature flag disabled for this ext
      - file size exceeds LOCAL_EXTRACT_MAX_MB
    Caller should fall back to LlamaParse / Docling in that case.
    """
    flags = {
        "docx": USE_LOCAL_EXTRACT_DOCX,
        "pptx": USE_LOCAL_EXTRACT_PPTX,
        "xlsx": USE_LOCAL_EXTRACT_XLSX,
    }
    if not flags.get(ext.lower(), False):
        return False
    if file_size_bytes and (file_size_bytes / (1024 * 1024)) > LOCAL_EXTRACT_MAX_MB:
        return False
    return True


# ============================================================
# DOCX (python-docx)
# ============================================================

def read_docx_as_markdown(filepath: str) -> Tuple[str, List[str]]:
    """Extract DOCX as markdown via python-docx (Thai-safe, fast).

    Preserves heading levels (Heading 1 → #, Heading 2 → ##, ...) and renders
    tables as GFM markdown. Speaker comments / footnotes / endnotes are not
    extracted (rare in practical PDB files).
    """
    warnings: List[str] = []
    try:
        from docx import Document
    except ImportError:
        return ("[docx extractor unavailable: install python-docx]", ["docx-dep-missing"])

    try:
        size_mb = os.path.getsize(filepath) / (1024 * 1024)
        if size_mb > LOCAL_EXTRACT_MAX_MB:
            warnings.append(f"large-file:{size_mb:.1f}MB")
    except OSError:
        pass

    try:
        with PARSE_LOCK:
            return _read_docx_unlocked(filepath, Document, warnings)
    except Exception as e:
        logger.error("DOCX local extract failed for %s: %s", filepath, e)
        return (f"[Extraction error: {type(e).__name__}: {str(e)[:200]}]", warnings + ["local-failed"])


def _read_docx_unlocked(filepath: str, Document, warnings: List[str]) -> Tuple[str, List[str]]:
    doc = Document(filepath)
    parts: List[str] = []
    for para in doc.paragraphs:
        txt = para.text.strip()
        if not txt:
            continue
        style = (para.style.name if para.style else "").strip()
        if style.startswith("Heading"):
            try:
                level = int(style.replace("Heading", "").strip())
            except (ValueError, AttributeError):
                level = 2
            parts.append(("#" * max(1, min(level, 6))) + " " + txt)
        else:
            parts.append(txt)
    for table in doc.tables:
        rows_md: List[str] = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip().replace("|", "\\|").replace("\n", " ") for c in row.cells]
            rows_md.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
        if rows_md:
            parts.append("\n".join(rows_md))
    text = "\n\n".join(parts)
    if not text.strip():
        return ("[No text content found in DOCX]", warnings + ["empty-docx"])
    return (text, warnings)


# ============================================================
# PPTX (python-pptx)
# ============================================================

def read_pptx_as_markdown(filepath: str) -> Tuple[str, List[str]]:
    """Extract PPTX as markdown via python-pptx — slide title + body + notes."""
    warnings: List[str] = []
    try:
        from pptx import Presentation
    except ImportError:
        return ("[pptx extractor unavailable: install python-pptx]", ["pptx-dep-missing"])

    try:
        size_mb = os.path.getsize(filepath) / (1024 * 1024)
        if size_mb > LOCAL_EXTRACT_MAX_MB:
            warnings.append(f"large-file:{size_mb:.1f}MB")
    except OSError:
        pass

    try:
        with PARSE_LOCK:
            prs = Presentation(filepath)
        sections: List[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
            except Exception:
                pass
            sections.append(f"## Slide {i}" + (f": {title}" if title else ""))
            body_lines: List[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    if hasattr(slide.shapes, "title") and shape == slide.shapes.title:
                        continue
                except Exception:
                    pass
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        body_lines.append(f"- {line}")
            if body_lines:
                sections.append("\n".join(body_lines))
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        sections.append(f"**Notes:** {notes}")
            except Exception:
                pass
        text = "\n\n".join(sections)
        if not text.strip():
            return ("[Empty presentation]", warnings + ["empty-pptx"])
        return (text, warnings)
    except Exception as e:
        logger.error("PPTX local extract failed for %s: %s", filepath, e)
        return (f"[Extraction error: {type(e).__name__}: {str(e)[:200]}]", warnings + ["local-failed"])


# ============================================================
# XLSX (openpyxl)
# ============================================================

def read_xlsx_as_markdown(filepath: str) -> Tuple[str, List[str]]:
    """Extract XLSX as markdown via openpyxl (read_only + data_only).

    HANDOFF Lesson 6: read_only=True streams rows to keep RAM bounded
    even on large workbooks; data_only=True resolves formulas to cached
    values (otherwise we'd get the formula string instead of the result).
    """
    warnings: List[str] = []
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ("[xlsx extractor unavailable: install openpyxl]", ["xlsx-dep-missing"])

    try:
        size_mb = os.path.getsize(filepath) / (1024 * 1024)
        if size_mb > LOCAL_EXTRACT_MAX_MB:
            warnings.append(f"large-file:{size_mb:.1f}MB")
    except OSError:
        pass

    # openpyxl does NOT need PARSE_LOCK — it doesn't share the lxml parser
    # (uses its own SAX-style reader in read_only mode).
    try:
        wb = load_workbook(filepath, data_only=True, read_only=True)
        sections: List[str] = []
        for sheet in wb.worksheets:
            sections.append(f"## Sheet: {sheet.title}")
            rows_md: List[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    safe = [c.replace("|", "\\|").replace("\n", " ") for c in cells]
                    rows_md.append("| " + " | ".join(safe) + " |")
            if rows_md:
                sections.append("\n".join(rows_md))
            else:
                sections.append("_(empty sheet)_")
                warnings.append(f"empty-sheet:{sheet.title}")
        wb.close()
        text = "\n\n".join(sections)
        if not text.strip():
            return ("[Empty workbook]", warnings + ["empty-xlsx"])
        return (text, warnings)
    except Exception as e:
        logger.error("XLSX local extract failed for %s: %s", filepath, e)
        return (f"[Extraction error: {type(e).__name__}: {str(e)[:200]}]", warnings + ["local-failed"])
