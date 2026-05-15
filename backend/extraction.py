"""
Text extraction — v10.0.0.

Pipeline:
1. Docling (IBM) for structured document understanding (if available)
2. PyPDF2 basic text extraction for PDF
3. pytesseract OCR fallback for image-only PDFs (if available)
4. Thai text post-processing to fix spacing issues
5. v10.0.0 — LlamaParse for PDF (opt-in via USE_LLAMAPARSE_FOR_PDF env)

Safety (v10.0.0):
  - PARSE_LOCK (threading.Lock) wraps python-docx + python-pptx calls
    to serialize lxml's module-level oxml_parser. Without this,
    concurrent worker jobs intermittently raise XMLSyntaxError
    (HANDOFF Lesson 2).

Reference: https://github.com/DS4SD/docling
"""
import os
import re
import logging

from .processors.safety import PARSE_LOCK

logger = logging.getLogger(__name__)


# v10.0.0 -- pre-compile regex used inside _postprocess_thai. Previously
# re.compile() ran once per call, recompiling the pattern on every PDF/image
# extraction. Module-level compile cuts the per-call overhead to ~zero.
_THAI_COMBINING_RE = re.compile(r'(\S) ([ัิ-ฺ็-๎])')
_SPACED_LATIN_RE = re.compile(r'(?:[A-Za-z0-9] ){3,}[A-Za-z0-9]')
_MULTI_SPACE_RE = re.compile(r'[^\S\n]{2,}')


# v9.3.3 — Strip lone surrogate code points before passing text to DB / hash / LLM.
#
# Why: PDF extraction (PyPDF2 / OCR / pdfplumber) sometimes emits text containing
# unpaired UTF-16 surrogates (U+D800–U+DFFF) from font encoding edge cases —
# e.g., embedded fonts using Private Use Area glyphs that map to surrogate
# halves. Python `str` accepts these (internal UCS-4) but they fail UTF-8 encode:
#   - aiosqlite cursor.execute → UnicodeEncodeError when SQLite stores text
#   - hashlib.sha256 .encode() → UnicodeEncodeError
#   - JSON serialize for OpenRouter / Drive push → UnicodeEncodeError
#
# Fix at this boundary covers all downstream writers in one place. Lossy: the
# offending positions are replaced with U+FFFD (replacement character). For
# text where this happens, the loss is rare characters that PDF couldn't map
# correctly anyway — better than hard-failing the entire upload.
#
# Encountered in production 2026-05-08: position 12562 + position 263 across
# different PDFs. Filed under DUP-005 (boundary sanitization).
def strip_surrogates(text: str) -> str:
    """Replace lone surrogates with U+FFFD so the result encodes safely as UTF-8.

    No-op for empty / clean text (encode-decode round-trip is fast in CPython
    when no replacement is needed).
    """
    if not text:
        return text
    return text.encode("utf-8", errors="replace").decode("utf-8", errors="replace")

# ─── Feature detection ───
_HAS_DOCLING = False
try:
    from docling.document_converter import DocumentConverter
    _HAS_DOCLING = True
    logger.info("Docling available — using advanced document understanding")
except ImportError:
    logger.warning("Docling not available — falling back to basic extraction")

_HAS_OCR = False
try:
    import pytesseract
    from pdf2image import convert_from_path
    _HAS_OCR = True
    logger.info("OCR available (pytesseract + pdf2image)")
except ImportError:
    logger.warning("OCR not available — PDF image extraction disabled")

# v9.0.0 — HEIC/HEIF support (iPhone photos)
# Without this, PIL.Image.open(.heic) raises UnidentifiedImageError
_HAS_HEIF = False
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
    _HAS_HEIF = True
    logger.info("HEIF/HEIC support enabled (pillow-heif)")
except ImportError:
    logger.warning("HEIC/HEIF not available — install pillow-heif for iPhone photo support")


def extract_text(filepath: str, filetype: str, progress_callback=None) -> str:
    """Extract text from a file (public API — sanitizes lone surrogates).

    Supported (v7.5.0):
      - pdf, docx (Docling structured + PyPDF2/python-docx fallbacks)
      - txt, md (encoding-tolerant)
      - png, jpg, jpeg, webp (Tesseract OCR — Thai+English)
      - csv (treated as txt)

    Returns extracted text. Error markers are wrapped in `[brackets]` so
    downstream `compute_content_hash()` skips hashing them (avoids
    false-positive matches between failed extractions).

    v9.3.3: result is filtered through strip_surrogates() so PDF extractions
    that emit lone UTF-16 halves never reach DB / hash / LLM downstream.

    v9.4.0: optional `progress_callback(step: str, pct: int|None)` — sync function
    เรียกระหว่าง extract เพื่อ report ขั้นตอนสด (TC-1 truthful: pct=None ถ้าไม่รู้)
    Default=None → backward compat ทุก existing caller (main.py, line_bot, mcp_tools, etc.)
    """
    return strip_surrogates(_extract_text_raw(filepath, filetype, progress_callback))


def _extract_text_raw(filepath: str, filetype: str, progress_callback=None) -> str:
    """Internal extraction without surrogate sanitization. See extract_text()."""
    try:
        # v10.0.0 — LlamaParse for PDF (opt-in, falls back to Docling/pypdf/OCR)
        if filetype == "pdf":
            llama_text = _try_llamaparse_pdf(filepath, progress_callback)
            if llama_text:
                return _postprocess_thai(llama_text)
            # else: fall through to Docling/pypdf/OCR chain below

        # v10.0.2 — HANDOFF Decision 4: DOCX/PPTX/XLSX → local extract FIRST
        # (skip Docling — local is 5-6x faster + better Thai accuracy).
        # Falls back to LlamaParse for DOCX if local fails.
        if filetype in ("docx", "pptx", "xlsx"):
            local_text = _try_local_office(filepath, filetype, progress_callback)
            if local_text and not local_text.startswith("["):
                return local_text
            # Local failed or skipped → for docx try LlamaParse as fallback
            if filetype == "docx" and _try_llamaparse_pdf is not None:
                try:
                    from .config import is_llamaparse_configured
                    if is_llamaparse_configured():
                        _safe_progress(progress_callback, "LlamaParse DOCX (fallback)", 50)
                        from .processors.llamaparse import parse_pdf
                        llama_text = parse_pdf(filepath, language="th")
                        if llama_text:
                            return _postprocess_thai(llama_text)
                except Exception as e:
                    logger.warning("LlamaParse DOCX fallback failed for %s: %s", filepath, e)
            # Final: return local result (marker or partial)
            return local_text or "[No text content found]"

        if filetype == "pdf" and _HAS_DOCLING:
            text = _extract_with_docling(filepath)
            if text and not text.startswith("["):
                return _postprocess_thai(text)
            # Docling returned nothing, try fallback
            return _extract_pdf_with_fallbacks(filepath, progress_callback)
        elif filetype == "pdf":
            return _extract_pdf_with_fallbacks(filepath, progress_callback)
        elif filetype in ("txt", "md", "csv",
                          # v9.0.0 — code files (text-based, encoding fallback ใช้ได้เลย)
                          "py", "js", "ts", "jsx", "tsx",          # Python / JS / TS
                          "css", "scss", "less", "sass",            # Stylesheets
                          "xml", "yaml", "yml", "toml", "ini",      # Config
                          "env", "conf", "cfg",                     # More config
                          "sh", "bash", "zsh", "bat", "ps1",        # Shell scripts
                          "sql",                                    # Database queries
                          "java", "kt", "swift",                    # Mobile
                          "c", "cpp", "h", "hpp", "cs",             # Compiled
                          "go", "rs", "rb", "php",                  # Other languages
                          "log", "tsv",                             # Data logs
                          "vue", "svelte"):                         # Components
            return _extract_txt(filepath)
        elif filetype in ("png", "jpg", "jpeg", "webp",
                          # v9.0.0 — extra image formats via PIL native + Tesseract OCR
                          "heic", "heif",   # iPhone photos (needs pillow-heif)
                          "gif", "bmp", "tiff", "tif"):  # PIL native support
            return _extract_image_ocr(filepath, progress_callback)
        elif filetype == "xlsx":
            return _extract_xlsx(filepath)
        elif filetype == "pptx":
            return _extract_pptx(filepath)
        elif filetype == "html":
            return _extract_html(filepath)
        elif filetype == "json":
            return _extract_json(filepath)
        elif filetype == "rtf":
            return _extract_rtf(filepath)
        else:
            # v9.0.0 — Phase B v2: ตรวจ AI multimodal formats (audio/video)
            # Return marker — caller (main.py upload) จะ route ไป ai_ingest.ingest_via_ai()
            from .ai_ingest import is_ai_format
            if is_ai_format(filetype):
                return f"[AI ingest needed: {filetype}]"
            return f"[Unsupported file type: {filetype}]"
    except Exception as e:
        logger.error(f"Extraction failed for {filepath}: {e}")
        # Try fallback
        try:
            if filetype == "pdf":
                return _extract_pdf_with_fallbacks(filepath, progress_callback)
            elif filetype == "docx":
                return _extract_docx_basic(filepath)
            elif filetype in ("png", "jpg", "jpeg", "webp",
                              "heic", "heif", "gif", "bmp", "tiff", "tif"):
                return _extract_image_ocr(filepath, progress_callback)
        except Exception:
            pass
        return f"[Extraction error: {str(e)}]"


def _try_local_office(filepath: str, filetype: str, progress_callback=None) -> str:
    """v10.0.2 — HANDOFF Decision 4: DOCX/PPTX/XLSX → local extract first.

    Returns extracted markdown text, or empty string if local is skipped
    (feature flag off or file too large) so caller can fall back. Returns
    a "[bracket-marker]" string if local was tried but produced no usable
    text — caller can decide whether to fallback further.

    Warnings list is logged + stored on a thread-local for caller to read
    (see _last_local_warnings).
    """
    try:
        from .processors.local import (
            should_use_local,
            read_docx_as_markdown,
            read_pptx_as_markdown,
            read_xlsx_as_markdown,
        )
    except ImportError:
        return ""

    try:
        size_bytes = os.path.getsize(filepath)
    except OSError:
        size_bytes = 0

    if not should_use_local(filetype, size_bytes):
        logger.info("Local extract skipped for %s (flag off or size > LOCAL_EXTRACT_MAX_MB)", filetype)
        return ""

    logger.info("Local extract START for %s: %s (%.2f MB)", filetype, os.path.basename(filepath), size_bytes / (1024 * 1024))
    _safe_progress(progress_callback, f"อ่าน {filetype.upper()} (local)", 35)
    if filetype == "docx":
        text, warnings = read_docx_as_markdown(filepath)
    elif filetype == "pptx":
        text, warnings = read_pptx_as_markdown(filepath)
    elif filetype == "xlsx":
        text, warnings = read_xlsx_as_markdown(filepath)
    else:
        return ""

    if warnings:
        logger.info("local %s warnings for %s: %s", filetype, os.path.basename(filepath), warnings)
        _LAST_LOCAL_WARNINGS.warnings = warnings  # noqa — thread-local stash
    if text and not text.startswith("["):
        logger.info("local %s: %d chars from %s", filetype, len(text), os.path.basename(filepath))
    return text


# Thread-local stash for warnings collected during the most recent local
# extract — read by upload_worker after extract_text returns. Avoids changing
# the public extract_text() signature (Pattern G partial implementation).
import threading as _threading
_LAST_LOCAL_WARNINGS = _threading.local()


def get_last_local_warnings() -> list:
    """Return + clear the warnings collected during the most recent extract_text() call."""
    w = getattr(_LAST_LOCAL_WARNINGS, "warnings", None) or []
    _LAST_LOCAL_WARNINGS.warnings = []
    return list(w)


def _try_llamaparse_pdf(filepath: str, progress_callback=None) -> str:
    """v10.0.0 -- attempt LlamaParse for PDF. Empty string -> caller should fallback.

    Behavior:
      * Returns extracted markdown text on success.
      * Returns "" (empty) if LlamaParse is not configured / not installed /
        raised an error -- caller falls back to Docling/pypdf/OCR.
      * Never raises -- failure is logged + caller continues.
    """
    try:
        from .config import is_llamaparse_configured
    except ImportError:
        return ""
    if not is_llamaparse_configured():
        return ""

    _safe_progress(progress_callback, "LlamaParse PDF", 25)
    try:
        from .processors.llamaparse import parse_pdf
        # v10.0.3 — forward progress_callback so poll loop reports "LlamaParse กำลังประมวลผล (Ns)"
        def _proxy(step: str, pct=None):
            _safe_progress(progress_callback, step, pct)
        text = parse_pdf(filepath, language="th", progress_callback=_proxy)
        if text and text.strip():
            logger.info("LlamaParse: %d chars from %s", len(text), os.path.basename(filepath))
            return text
        logger.warning("LlamaParse returned empty for %s -- will fallback", os.path.basename(filepath))
        return ""
    except Exception as e:
        logger.warning("LlamaParse failed for %s: %s -- will fallback", os.path.basename(filepath), e)
        return ""


def _safe_progress(progress_callback, step: str, pct=None) -> None:
    """v9.4.0 helper — เรียก progress_callback แบบปลอดภัย (catch all errors).

    progress_callback อาจเป็น None (callers เก่า) หรือ raise exception ระหว่าง update DB.
    Wrap ที่นี่เพื่อ extract path ไม่พังเพราะ progress reporting fail.
    """
    if progress_callback is None:
        return
    try:
        progress_callback(step, pct)
    except Exception as e:
        logger.debug(f"progress_callback raised (non-fatal): {e}")


def _extract_image_ocr(filepath: str, progress_callback=None) -> str:
    """OCR text from png/jpg/jpeg/webp via pytesseract (Thai + English).

    v7.5.0: รองรับ image upload จริง (ก่อนหน้านี้ png/jpg ถูกตั้งใน allowed_types
    แต่ extraction.py reject กลับเป็น "[Unsupported file type]" → orphan ใน DB)

    v9.4.0: optional progress_callback — report "OCR รูปภาพ" before tesseract call

    Returns:
      - Extracted text (post-processed for Thai) ถ้าเจอ
      - "[Image: no text detected]" ถ้า OCR สำเร็จแต่ไม่มี text
      - "[Image: OCR not available]" ถ้า pytesseract / tesseract binary หาย
      - "[OCR error: ...]" ถ้า exception อื่น
    """
    if not _HAS_OCR:
        return "[Image: OCR not available]"
    try:
        _safe_progress(progress_callback, "เปิดรูปภาพ", 30)
        from PIL import Image
        img = Image.open(filepath)
        # Convert mode if needed (RGBA / palette → RGB for tesseract)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        _safe_progress(progress_callback, "OCR รูปภาพ", 60)
        text = pytesseract.image_to_string(img, lang="tha+eng")
        if text and text.strip():
            logger.info(f"Image OCR: extracted {len(text)} chars from {os.path.basename(filepath)}")
            _safe_progress(progress_callback, "ประมวลผลข้อความ Thai", 85)
            return _postprocess_thai(text.strip())
        return "[Image: no text detected]"
    except Exception as e:
        logger.error(f"Image OCR failed for {filepath}: {e}")
        return f"[OCR error: {str(e)}]"


def _extract_pdf_with_fallbacks(filepath: str, progress_callback=None) -> str:
    """PDF extraction with full fallback chain: encrypted check → PyPDF2 → OCR.

    v7.5.0: detects encrypted PDF early so caller can flag extraction_status.
    v9.4.0: optional progress_callback — report each phase + page-level for OCR
    """
    # v7.5.0 — short-circuit for encrypted PDFs (specific marker for downstream)
    _safe_progress(progress_callback, "ตรวจไฟล์ PDF", 25)
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        if reader.is_encrypted:
            logger.info(f"PDF encrypted: {os.path.basename(filepath)}")
            return "[PDF encrypted: password-protected — unlock before re-uploading]"
    except Exception:
        # If we can't even open the file, fall through to normal error handling
        pass

    # Step 1: Try PyPDF2 text extraction
    _safe_progress(progress_callback, "อ่าน PDF (text layer)", 35)
    text = _extract_pdf_basic(filepath, progress_callback)

    if text and not text.startswith("[No text"):
        return _postprocess_thai(text)

    # Step 2: Try OCR if text extraction failed
    if _HAS_OCR:
        logger.info(f"PyPDF2 found no text — trying OCR for {os.path.basename(filepath)}")
        _safe_progress(progress_callback, "PDF เป็นรูปสแกน — เริ่ม OCR", 50)
        ocr_text = _extract_pdf_ocr(filepath, progress_callback)
        if ocr_text and not ocr_text.startswith("["):
            return _postprocess_thai(ocr_text)

    # Step 3: v10.0.1 — Gemini PDF fallback (handles image-only PDFs even
    # when Tesseract is unavailable; supports up to 1000 pages / 20MB inline).
    try:
        from .ai_ingest import is_available as _ai_available, extract_pdf_via_gemini_sync
        if _ai_available():
            logger.info(f"OCR unavailable / failed — trying Gemini PDF for {os.path.basename(filepath)}")
            _safe_progress(progress_callback, "ส่ง Gemini อ่าน PDF", 70)
            gem_text = extract_pdf_via_gemini_sync(filepath)
            if gem_text and not gem_text.startswith("["):
                return _postprocess_thai(gem_text)
            # If Gemini returned a marker, surface it (more informative than generic)
            if gem_text:
                return gem_text
    except Exception as e:
        logger.warning("Gemini PDF fallback errored for %s: %s", filepath, e)

    return "[No text content found in PDF — file may be image-only and OCR is not available]"


def classify_extraction_status(text: str) -> str:
    """v7.5.0 — derive extraction_status from extracted_text content.

    Maps the [bracket-marker] convention used by extract_text() back into the
    short status enum stored in the DB. Single source of truth for status
    classification — used by upload, reprocess, and migration paths.

    Returns one of: ok | empty | encrypted | ocr_failed | unsupported
    """
    if not text or not text.strip():
        return "empty"
    if not text.startswith("["):
        return "ok"
    lower = text.lower()
    # Empty markers
    if lower.startswith("[empty"):
        return "empty"
    if "encrypted" in lower or "password-protected" in lower:
        return "encrypted"
    # AI ingest configuration / unsupported (caller hasn't set up AI properly)
    if "ai ingest not configured" in lower or "ai ingest unsupported" in lower:
        return "unsupported"
    if "unsupported file type" in lower:
        return "unsupported"
    # AI ingest / vision errors → ocr_failed (retryable)
    if (
        "ai ingest error" in lower
        or "ai vision error" in lower
        or "ai image: no description" in lower
        or "ai audio: no transcription" in lower
        or "ai video: no analysis" in lower
        or "image: ocr not available" in lower
        or "image: no text detected" in lower
    ):
        return "ocr_failed"
    if "no text detected" in lower or "no text content" in lower or "ocr found no text" in lower:
        return "ocr_failed"
    if "ocr error" in lower or "extraction error" in lower:
        return "ocr_failed"
    # Extractor missing dependencies (xlsx/pptx/rtf/html unavailable)
    if "extractor unavailable" in lower:
        return "unsupported"
    return "ok"  # safe default — unknown bracket marker shouldn't block


def _extract_with_docling(filepath: str) -> str:
    """
    Extract text using Docling — produces structured Markdown.
    Preserves headings, tables, lists, and document hierarchy.
    """
    logger.info(f"Docling: processing {os.path.basename(filepath)}")
    converter = DocumentConverter()
    result = converter.convert(filepath)

    # Export to Markdown (preserves structure)
    md_text = result.document.export_to_markdown()

    if md_text and md_text.strip():
        logger.info(f"Docling: extracted {len(md_text)} chars from {os.path.basename(filepath)}")
        return md_text.strip()

    return "[No text content found]"


# ─── FALLBACK EXTRACTORS ───

def _extract_pdf_basic(filepath: str, progress_callback=None) -> str:
    """Fallback PDF extraction using PyPDF2.

    v9.4.0: optional progress_callback report per-page progress.
    PDF text layer extraction = fast (sub-second per page) so report ทุก 5 pages
    """
    from PyPDF2 import PdfReader
    reader = PdfReader(filepath)
    pages = []
    total = len(reader.pages)
    for i, page in enumerate(reader.pages):
        if progress_callback and (i % 5 == 0 or i == total - 1):
            # Report 30-60% range during PyPDF2 phase (after PDF check at 25%, before OCR at 50%)
            pct = 30 + int(((i + 1) / total) * 30) if total > 0 else 35
            _safe_progress(progress_callback, f"อ่าน PDF หน้า {i+1}/{total}", pct)
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"--- Page {i+1} ---\n{text.strip()}")
    return "\n\n".join(pages) if pages else "[No text content found in PDF]"


def _extract_pdf_ocr(filepath: str, progress_callback=None) -> str:
    """OCR extraction using pytesseract for image-only PDFs.

    v9.4.0: optional progress_callback report per-page OCR progress
    OCR = slow (1-3s per page) so report ทุก page
    """
    try:
        _safe_progress(progress_callback, "แปลง PDF → รูปภาพ", 55)
        images = convert_from_path(filepath, dpi=200, first_page=1, last_page=20)
        pages = []
        total = len(images)
        for i, img in enumerate(images):
            if progress_callback:
                # OCR phase = 60-90% range (after image conversion at 55%, before save at 95%)
                pct = 60 + int(((i + 1) / total) * 30) if total > 0 else 75
                _safe_progress(progress_callback, f"OCR หน้า {i+1}/{total}", pct)
            # Use Thai + English language for OCR
            text = pytesseract.image_to_string(img, lang='tha+eng')
            if text and text.strip():
                pages.append(f"--- Page {i+1} (OCR) ---\n{text.strip()}")

        if pages:
            logger.info(f"OCR: extracted {sum(len(p) for p in pages)} chars from {os.path.basename(filepath)}")
            return "\n\n".join(pages)

        return "[OCR found no text in PDF images]"
    except Exception as e:
        logger.error(f"OCR failed for {filepath}: {e}")
        return f"[OCR error: {str(e)}]"


def _extract_docx_basic(filepath: str) -> str:
    """Fallback DOCX extraction using python-docx.

    v10.0.0: serialized via PARSE_LOCK -- python-docx uses lxml's
    module-level oxml_parser which is NOT thread-safe (HANDOFF Lesson 2).
    """
    from docx import Document
    with PARSE_LOCK:
        return _extract_docx_basic_unlocked(filepath, Document)


def _extract_docx_basic_unlocked(filepath: str, Document) -> str:
    doc = Document(filepath)
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            # Preserve heading structure
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading ', '').strip()
                try:
                    hashes = '#' * int(level)
                except ValueError:
                    hashes = '##'
                paragraphs.append(f"{hashes} {para.text.strip()}")
            else:
                paragraphs.append(para.text.strip())
    # Extract tables as markdown
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        if rows:
            paragraphs.append("\n".join(rows))
    return "\n\n".join(paragraphs) if paragraphs else "[No text content found in DOCX]"


def _extract_txt(filepath: str) -> str:
    """Extract text from TXT/MD files with encoding fallback."""
    encodings = ['utf-8', 'utf-8-sig', 'cp874', 'tis-620', 'latin-1']
    for enc in encodings:
        try:
            with open(filepath, 'r', encoding=enc) as f:
                text = f.read()
            return text.strip() if text.strip() else "[Empty file]"
        except (UnicodeDecodeError, UnicodeError):
            continue
    return "[Could not decode file]"


# ─── TEXT POST-PROCESSING (LLM-powered) ───

async def cleanup_extracted_text(text: str, filename: str) -> str:
    """Use LLM to clean up broken PDF text extraction.
    
    PDF extractors often produce broken Thai text with:
    - Character spacing: ผ ่า น เ ข ้า ร อ บ
    - Private Use chars: \\uf70a instead of ่
    - Spaced Latin: S u b s c r i p t i o n
    
    LLM understands language context and fixes everything in one pass.
    """
    from .llm import call_llm_pro
    
    if not text or len(text) < 50:
        return text
    
    # Check if text needs cleanup
    space_ratio = text.count(' ') / max(len(text), 1)
    has_private_use = any(0xF700 <= ord(c) <= 0xF7FF for c in text[:1000])
    
    if space_ratio < 0.25 and not has_private_use:
        logger.info(f"Text looks clean (space ratio {space_ratio:.1%}), skipping LLM cleanup")
        return text
    
    logger.info(f"LLM cleanup: {filename} ({len(text)} chars, space ratio {space_ratio:.1%})")
    
    # Process in chunks if text is very long
    max_chunk = 6000
    if len(text) > max_chunk:
        chunks = []
        for i in range(0, len(text), max_chunk):
            chunk = text[i:i + max_chunk]
            cleaned = await _llm_fix_chunk(chunk, filename)
            chunks.append(cleaned)
        result = "\n".join(chunks)
    else:
        result = await _llm_fix_chunk(text, filename)
    
    logger.info(f"LLM cleanup done: {len(text)} → {len(result)} chars")
    return result


async def _llm_fix_chunk(text: str, filename: str) -> str:
    """Send a chunk of broken text to LLM for cleanup."""
    from .llm import call_llm_pro
    
    system_prompt = """คุณเป็นผู้เชี่ยวชาญในการแก้ไขข้อความที่ extract จาก PDF

กฎ:
1. แก้ไขข้อความที่ตัวอักษรถูกเว้นวรรคผิด เช่น "ผ ่า น เ ข ้า" → "ผ่านเข้า"
2. แก้ไขตัวอักษร Unicode Private Use Area (\\uf70a, \\uf70b ฯลฯ) ให้เป็นสระ/วรรณยุกต์ไทยที่ถูกต้อง
3. แก้ไขตัวอักษรภาษาอังกฤษที่ถูกเว้น เช่น "S u b s c r i p t i o n" → "Subscription"  
4. รักษาเนื้อหาและความหมายเดิมทั้งหมด ห้ามเพิ่ม/ลบ/แต่งเนื้อหา
5. รักษา line breaks และ formatting เดิม (เช่น --- Page 1 ---)
6. ถ้าข้อความปกติอยู่แล้ว ให้คืนตามเดิม
7. ตอบเฉพาะข้อความที่แก้ไขแล้วเท่านั้น ไม่ต้องมีคำอธิบาย"""

    user_prompt = f"แก้ไขข้อความที่ extract จาก {filename}:\n\n{text}"
    
    try:
        result = await call_llm_pro(system_prompt, user_prompt, temperature=0.1, max_tokens=16384)
        return result.strip()
    except Exception as e:
        logger.error(f"LLM cleanup failed for {filename}: {e}")
        return text


def _postprocess_thai(text: str) -> str:
    """Legacy sync wrapper — basic regex cleanup only.
    For full cleanup, use async cleanup_extracted_text() instead.
    """
    if not text:
        return text
    # Basic: just remove spaces before Thai combining characters
    prev = None
    result = text
    while prev != result:
        prev = result
        result = _THAI_COMBINING_RE.sub(r'\1\2', result)
    # Collapse spaced Latin
    def fix_latin(m):
        return m.group(0).replace(' ', '')
    result = _SPACED_LATIN_RE.sub(fix_latin, result)
    result = _MULTI_SPACE_RE.sub(' ', result)
    return result



# ─── v7.5.0 — Phase 3: more format extractors ────────────────────────


def _extract_xlsx(filepath: str) -> str:
    """Extract text from xlsx (Excel) — flatten all sheets to markdown tables.

    Multi-sheet decision (Q): flatten ทุก sheet เป็นไฟล์เดียว (per user choice).
    Uses data_only=True so formulas resolve to cached values (not formula text).
    read_only=True for memory efficiency on big workbooks.
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[xlsx extractor unavailable: install openpyxl]"
    try:
        wb = load_workbook(filepath, data_only=True, read_only=True)
        sections = []
        for sheet in wb.worksheets:
            sections.append(f"## Sheet: {sheet.title}")
            rows_md = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows_md.append("| " + " | ".join(cells) + " |")
            if rows_md:
                sections.append("\n".join(rows_md))
            else:
                sections.append("_(empty sheet)_")
        return "\n\n".join(sections) if sections else "[Empty workbook]"
    except Exception as e:
        logger.error(f"xlsx extraction failed for {filepath}: {e}")
        return f"[Extraction error: {str(e)}]"


def _extract_pptx(filepath: str) -> str:
    """Extract text + speaker notes from pptx (PowerPoint).

    Each slide -> "## Slide N: <title>" section with body bullets + notes.

    v10.0.0: serialized via PARSE_LOCK (python-pptx uses lxml).
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "[pptx extractor unavailable: install python-pptx]"
    try:
        with PARSE_LOCK:
            prs = Presentation(filepath)
        sections = []
        for i, slide in enumerate(prs.slides, start=1):
            # Title
            title = ""
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
            except Exception:
                pass
            sections.append(f"## Slide {i}" + (f": {title}" if title else ""))
            # Body text (skip title shape — already captured)
            body_lines = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if hasattr(slide.shapes, "title") and shape == slide.shapes.title:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        body_lines.append(f"- {line}")
            if body_lines:
                sections.append("\n".join(body_lines))
            # Speaker notes
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        sections.append(f"**Notes:** {notes}")
            except Exception:
                pass
        return "\n\n".join(sections) if sections else "[Empty presentation]"
    except Exception as e:
        logger.error(f"pptx extraction failed for {filepath}: {e}")
        return f"[Extraction error: {str(e)}]"


def _extract_html(filepath: str) -> str:
    """Extract plain text from HTML — STRIPS <script> and <style> tags.

    Security: HTML may contain <script>alert(1)</script> XSS payloads. Even
    though we don't render the output to a browser, we DO send it to LLM
    summary + vector index, so injected text could pollute downstream context.
    BeautifulSoup .decompose() removes the entire tag tree.
    """
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return "[html extractor unavailable: install beautifulsoup4]"
    try:
        # Try common encodings
        content = None
        for enc in ("utf-8", "utf-8-sig", "cp874", "tis-620", "latin-1"):
            try:
                with open(filepath, "r", encoding=enc) as f:
                    content = f.read()
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
        if content is None:
            return "[Could not decode HTML file]"
        soup = BeautifulSoup(content, "html.parser")
        # Security: remove script/style/iframe (XSS + side-channel)
        for tag in soup(["script", "style", "iframe", "object", "embed", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        # Collapse excessive whitespace
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text if text.strip() else "[HTML had no extractable text]"
    except Exception as e:
        logger.error(f"html extraction failed for {filepath}: {e}")
        return f"[Extraction error: {str(e)}]"


def _extract_json(filepath: str) -> str:
    """Pretty-print JSON for readable extracted_text.

    Falls back to raw text if JSON malformed (some "JSON" files are JSON5
    or have trailing commas — still useful as text input).
    """
    try:
        import json as _json
        for enc in ("utf-8", "utf-8-sig"):
            try:
                with open(filepath, "r", encoding=enc) as f:
                    raw = f.read()
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
        else:
            return "[Could not decode JSON file]"
        try:
            data = _json.loads(raw)
            return _json.dumps(data, indent=2, ensure_ascii=False)
        except _json.JSONDecodeError:
            # Malformed JSON — return raw text (still useful)
            return raw if raw.strip() else "[Empty JSON file]"
    except Exception as e:
        logger.error(f"json extraction failed for {filepath}: {e}")
        return f"[Extraction error: {str(e)}]"


def _extract_rtf(filepath: str) -> str:
    """Strip RTF control codes → plain text via striprtf."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        return "[rtf extractor unavailable: install striprtf]"
    try:
        for enc in ("utf-8", "latin-1", "cp1252"):
            try:
                with open(filepath, "r", encoding=enc) as f:
                    content = f.read()
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
        else:
            return "[Could not decode RTF file]"
        text = rtf_to_text(content, errors="ignore")
        return text.strip() if text.strip() else "[RTF had no extractable text]"
    except Exception as e:
        logger.error(f"rtf extraction failed for {filepath}: {e}")
        return f"[Extraction error: {str(e)}]"
